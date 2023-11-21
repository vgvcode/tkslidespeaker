from pptx_tools import utils  # needs python-pptx-interface installed
from pptx import Presentation
from pptx.util import Inches
import shutil
import os
import sys
import config as cfg
import platform

def create_images_with_notes_presentation(images_folder, notes_folder, output_pptx):
    n = count_files(images_folder)  # Replace 'n' with the number of images you have
    print(f"Creating a presentation with {n} image slides and notes...")

    prs = Presentation()

    for i in range(1, n + 1):  # Replace 'n' with the number of images you have
        slide_layout = prs.slide_layouts[5]  # Blank slide layout

        slide = prs.slides.add_slide(slide_layout)

        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)  # Adjust the width as needed old width 9
        height = Inches(6)  # Adjust the height as needed

        img_path = f"{images_folder}/Slide{i}.PNG"  # Replace 'PNG' with the actual image file format

        pic = slide.shapes.add_picture(img_path, left, top, width, height)

        # Load and add notes from a text file
        notes_path = f"{notes_folder}/Notes-Slide{i}.txt"
        with open(notes_path, 'r', encoding='utf-8') as file:
            notes = file.read()

        # Add the notes to the slide
        slide.notes_slide.notes_text_frame.text = notes

    prs.save(output_pptx)
    return n

def extract_and_save_notes(input_pptx, output_folder):
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    prs = Presentation(input_pptx)

    for slide_number, slide in enumerate(prs.slides, 1):
        notes_slide = slide.notes_slide
        notes = notes_slide.notes_text_frame.text
        slide_notes_file = f"{output_folder}/Notes-Slide{slide_number}.txt"
        
        with open(slide_notes_file, 'w', encoding='utf-8') as notes_file:
            notes_file.write(notes)

    return len(prs.slides)

def count_files(folder_path):
    file_count = 0
    for filename in os.listdir(folder_path):
        if os.path.isfile(os.path.join(folder_path, filename)):
            file_count += 1
    #print(f"Total number of files in the folder: {file_count}")
    return file_count

def convert(preso):
    inputFilePath = os.path.join(cfg.stagingFolder, preso) #~/../staging/ai-teacher.pptx

    inputFilePathElements = os.path.split(inputFilePath)
    fileNameWithExt = inputFilePathElements[1] #ai-teacher.pptx
    fileElements = fileNameWithExt.split(".")
    fileNameWithoutExt = fileElements[0]    #ai-teacher
    ext = fileElements[1] #pptx

    presentationBaseFolder = os.path.join(cfg.tmpFolder, "_" + fileNameWithoutExt) #~/../tmp/_ai-teacher
    imagesFolder = os.path.join(presentationBaseFolder, "images") #~/../tmp/_ai-teacher/images
    notesFolder = os.path.join(presentationBaseFolder, "notes") #~/../tmp/_ai-teacher/notes
    outputFilePath = os.path.join(cfg.tmpFolder, fileNameWithExt) #~/../tmp/ai-teacher.pptx

    #print("Base folder: " + cfg.tmpFolder) 
    #print("Presentation base folder: " + presentationBaseFolder)
    #print("Images folder: " + imagesFolder)
    #print("Notes folder: " + notesFolder)
    #print("Input file: " + inputFilePath)
    #print("Output file: " + outputFilePath)

    plat = platform.system()

    #do this only on windows platform
    if plat == "Windows":
        #extract all images one per slide
        if not os.path.exists(imagesFolder):
            os.makedirs(imagesFolder)

        print("Extracting all images one per slide...")
        utils.save_pptx_as_png(imagesFolder, inputFilePath, overwrite_folder=True)  # additional optional parameter overwrite_folder

        #extract all notes one per slide
        print("Extracting all notes one per slide...")
        n = extract_and_save_notes(inputFilePath, notesFolder)

        #create a new pptx with images and notes
        print("Creating a new pptx with images and notes...")
        create_images_with_notes_presentation(imagesFolder, notesFolder, outputFilePath)

        print("Cleaning up...")

        #delete presentation base folder
        shutil.rmtree(presentationBaseFolder)
    else:
        #if platform is not windows, just count the number of slides
        #and copy the input file to the output file
        print("Platform is not Windows. Skipping file conversion")
        prs = Presentation(inputFilePath)
        n = len(prs.slides)
        shutil.copyfile(inputFilePath, outputFilePath)


    return (outputFilePath, n)